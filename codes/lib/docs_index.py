"""Extract text from every shared/ document (pptx/pdf/docx/xlsx), chunk it, embed it via
Ollama, and cache the result so it's only recomputed when a source file actually changes."""

from pathlib import Path

from .common import Ollama, cache_dir, notes_dir, read_json, subject_dir, write_json

DOC_EXTS = {".pptx", ".pdf", ".docx", ".xlsx"}


def _extract_pptx(path: Path) -> list[tuple[str, str]]:
    from pptx import Presentation

    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in row.cells))
        if texts:
            out.append((f"slide {i}", "\n".join(texts)))
    return out


def _extract_pdf(path: Path) -> list[tuple[str, str]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        reader.decrypt("")  # most "encrypted" PDFs from course portals are just owner-locked, no real password
    out = []
    for i, page in enumerate(reader.pages, start=1):
        t = (page.extract_text() or "").strip()
        if t:
            out.append((f"page {i}", t))
    return out


def _extract_docx(path: Path) -> list[tuple[str, str]]:
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells))
    text = "\n".join(parts)
    return [("document", text)] if text.strip() else []


def _extract_xlsx(path: Path) -> list[tuple[str, str]]:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append(" | ".join("" if c is None else str(c) for c in row))
        if rows:
            out.append((f"sheet '{ws.title}'", "\n".join(rows)))
    return out


EXTRACTORS = {".pptx": _extract_pptx, ".pdf": _extract_pdf, ".docx": _extract_docx, ".xlsx": _extract_xlsx}


def _chunk_text(text: str, chunk_words: int) -> list[str]:
    words = text.split()
    return [" ".join(words[i:i + chunk_words]) for i in range(0, len(words), chunk_words)] or []


def build_or_update_index(subject: str, config: dict, force: bool = False) -> list[dict]:
    """Returns the full list of {source, section, text, embedding} chunks for a subject,
    rebuilding only entries for files that are new or whose mtime changed."""
    index_path = cache_dir(subject) / "doc_index.json"
    cached = read_json(index_path, default={"files": {}, "chunks": []})
    if force:
        cached = {"files": {}, "chunks": []}

    shared_dir = subject_dir(subject) / "shared"
    doc_paths = [p for p in shared_dir.glob("*") if p.suffix.lower() in DOC_EXTS]

    ollama = Ollama(config)
    chunk_words = config["retrieval"]["chunk_words"]

    known_files = cached["files"]
    all_chunks = [c for c in cached["chunks"] if c["source"] in {p.name for p in doc_paths}]
    changed = False

    for path in doc_paths:
        mtime = path.stat().st_mtime
        if known_files.get(path.name) == mtime:
            continue  # unchanged, keep existing chunks for this file
        changed = True
        print(f"[docs_index] indexing {path.name}")
        all_chunks = [c for c in all_chunks if c["source"] != path.name]
        extractor = EXTRACTORS[path.suffix.lower()]
        try:
            sections = extractor(path)
        except Exception as e:
            print(f"[docs_index] failed to extract {path.name}: {e}")
            sections = []
        for section_label, text in sections:
            for chunk in _chunk_text(text, chunk_words):
                if not chunk.strip():
                    continue
                embedding = ollama.embed(chunk)
                all_chunks.append({
                    "source": path.name,
                    "section": section_label,
                    "text": chunk,
                    "embedding": embedding,
                })
        known_files[path.name] = mtime

    if changed:
        write_json(index_path, {"files": known_files, "chunks": all_chunks})
        print(f"[docs_index] {subject}: {len(all_chunks)} chunks across {len(known_files)} files")

    return all_chunks
